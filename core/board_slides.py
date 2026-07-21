"""
core/board_slides.py — PT Drift Tracker board-slide export (.pptx).

Replaces app.py's original "Export as Board Slide" button, which shelled
out to a Node.js script (generate_pt_slide.js) via subprocess — that
script was never actually included anywhere in the project (checked: no
.js file exists in the repo), so the original button would have failed if
clicked even in the demo, regardless of whether Node was installed. This
is a from-scratch, pure-Python replacement using python-pptx, so it works
without any external toolchain and without depending on a file the repo
never shipped.

Draws on the exact same data markets_page.py's PT Drift Tracker tab
already computes and displays (see _render_pt_drift and its "Analyst PT
direction" table) — same 8-quarter pt_history, same first-vs-last-PT
drift-direction math — so the exported slide always matches what's on
screen. The Revision Momentum line (see core/risk_scorecard.py) is passed
in rather than recomputed here, so the slide and the IR Risk Dashboard
never say two different things about the same 8-quarter history.

"Midnight Executive" palette (dark navy / ice blue / white) — matches the
app's own dark-surface Markets/Risk cards rather than a generic default.
"""

from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.util import Inches, Pt

_NAVY = RGBColor(0x1E, 0x27, 0x61)
_NAVY_LIGHT = RGBColor(0x2A, 0x35, 0x78)
_ICE = RGBColor(0xCA, 0xDC, 0xFC)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GREEN = RGBColor(0x4A, 0xDE, 0x80)
_RED = RGBColor(0xF8, 0x71, 0x71)
_MUTED = RGBColor(0x94, 0xA3, 0xB8)
_GRID = RGBColor(0x33, 0x3F, 0x66)

_LINE_COLORS = [
    RGBColor(0x93, 0xC5, 0xFD), RGBColor(0x4A, 0xDE, 0x80), RGBColor(0xFC, 0xD3, 0x4D),
    RGBColor(0xF9, 0x71, 0x67), RGBColor(0xA5, 0xA5, 0xA5),
]
_STOCK_COLOR = _WHITE
# The stock-price line is visually set apart from every firm's PT line via
# THREE independent cues (color, dash, marker shape) rather than color alone
# — a colorblind-safe / grayscale-safe habit, and it also sidesteps the
# original bug where this collided with a red firm color at n=4 firms.


def _drift_rows(pt_hist, last_price):
    """Same first-vs-last-PT drift math as markets_page.py's
    _render_pt_drift "Analyst PT direction" table — kept in sync
    deliberately so the slide never disagrees with the on-screen table."""
    labels = pt_hist.get("labels", [])
    by_firm = pt_hist.get("by_firm", {})
    rows = []
    for firm, pts in by_firm.items():
        # aligned to a shared date axis: use the firm's own valid points, not pts[-1] (which is None
        # unless the firm revised on the very last date in the series).
        valid = [(labels[i], p) for i, p in enumerate(pts) if p is not None and i < len(labels)]
        active = len(valid) >= 1
        if len(valid) >= 2 and valid[0][1]:
            last_pt = valid[-1][1]
            chg_pct = round((last_pt - valid[0][1]) / valid[0][1] * 100, 1)
        else:
            last_pt, chg_pct = (valid[-1][1] if valid else None), None
        upside = round((last_pt - last_price) / last_price * 100, 1) if last_pt and last_price else None
        rows.append({"firm": firm, "pt": last_pt, "chg_pct": chg_pct, "upside": upside, "active": active})
    return rows


def _textbox(slide, left, top, width, height, lines, align_center=False):
    """lines: list of (text, size_pt, bold, color) tuples, one per paragraph."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align_center:
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def _stat_callout(slide, left, top, width, label, value, sub, value_color):
    _textbox(slide, left, top, width, Inches(1.15), [
        (label, 10, True, _MUTED),
        (value, 24, True, value_color),
    ] + ([(sub, 11, False, _MUTED)] if sub else []))


def _pt_drift_placeholder_slide(client_name, ticker):
    """Honest one-slide fallback when there's no real multi-period PT drift yet: shows the current
    PTs we DO hold (from the registry) and states that the drift history accrues as PTs are logged.
    Never draws an invented 8-quarter chart."""
    from config.client_config import analysts_with_pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _NAVY

    _textbox(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.55), [
        (f"{client_name} ({ticker}) — Price Target Drift Tracker", 26, True, _WHITE),
    ])
    _textbox(slide, Inches(0.5), Inches(0.92), Inches(12.3), Inches(0.35), [
        (f"Generated {datetime.now().strftime('%b %d, %Y')}", 11, False, _MUTED),
    ])

    lines = [("Current price targets", 16, True, _WHITE)]
    aps = sorted(analysts_with_pt(), key=lambda a: -(a.get("pt") or 0))
    if aps:
        for a in aps:
            rating = a.get("rating") or "no rating logged"
            lines.append((f"{a.get('firm')} — {a.get('name')}:   ${a.get('pt'):.2f}   ({rating})", 13, False, _ICE))
    else:
        lines.append(("No analyst price targets logged yet.", 13, False, _MUTED))
    lines.append((" ", 8, False, _MUTED))
    lines.append(("PT drift history is not yet tracked — it accumulates as each analyst's PT is logged "
                  "per quarter (Earnings → Model Intake). A multi-quarter drift chart appears here once "
                  "at least two periods of PTs are on file.", 12, False, _MUTED))
    _textbox(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.0), lines)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_pt_drift_slide(client_name, ticker, seed, revision_momentum=None):
    """Returns pptx bytes for a one-slide PT Drift Tracker board slide,
    built from `seed` (data.seed.consensus_estimates's dict — the same one
    markets_page.py's PT Drift Tracker renders on screen) and an optional
    revision_momentum dict {"headline": str, "detail": str, "status": str}
    from core/risk_scorecard.py, so the two never disagree."""
    pt_hist = seed.get("pt_history", {})
    labels = pt_hist.get("labels", [])
    stock_prices = pt_hist.get("stock_prices", [])
    by_firm = pt_hist.get("by_firm", {})
    last_price = stock_prices[-1] if stock_prices else 0
    if not last_price:
        # pt_history (rating-action feed) carries no price series, so pull the current price for the
        # LAST PRICE callout + upside math rather than showing $0.00.
        try:
            from core import market_data
            snap = market_data.get_snapshot(ticker)
            last_price = (snap or {}).get("last_price") or 0
        except Exception:
            last_price = 0

    # No real multi-date drift on file → render an honest current-PTs slide, never a fabricated
    # chart. (pt_history is built from the analyst rating-action feed by core.consensus; it's empty
    # until ≥2 revision dates exist.)
    if not by_firm or len(labels) < 2:
        return _pt_drift_placeholder_slide(client_name, ticker)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _NAVY

    _textbox(slide, Inches(0.5), Inches(0.35), Inches(8.0), Inches(0.55), [
        (f"{client_name} ({ticker}) — Price Target Drift Tracker", 26, True, _WHITE),
    ])
    _textbox(slide, Inches(0.5), Inches(0.92), Inches(8.0), Inches(0.35), [
        (f"Sell-side PT history (analyst rating-action feed) · generated {datetime.now().strftime('%b %d, %Y')}", 11, False, _MUTED),
    ])

    # Native, editable line chart — left ~62% of the slide.
    chart_data = CategoryChartData()
    chart_data.categories = labels
    for firm, pts in by_firm.items():
        chart_data.add_series(firm, tuple(pts))
    if stock_prices:
        chart_data.add_series(f"{ticker} Stock Price", tuple(stock_prices))

    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(0.5), Inches(1.45), Inches(7.9), Inches(5.55), chart_data,
    )
    chart = gframe.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = _WHITE
    plot = chart.plots[0]
    n_firms = len(by_firm)
    for i, series in enumerate(plot.series):
        is_stock = i == n_firms
        color = _STOCK_COLOR if is_stock else _LINE_COLORS[i % len(_LINE_COLORS)]
        series.format.line.color.rgb = color
        series.format.line.width = Pt(2.75 if is_stock else 2.25)
        if is_stock:
            # Distinct from every firm PT line on THREE cues at once (color,
            # dash, marker) — not relying on color alone to read as separate.
            series.format.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            series.marker.style = XL_MARKER_STYLE.DIAMOND
        else:
            series.marker.style = XL_MARKER_STYLE.CIRCLE
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = color
        series.smooth = False
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.color.rgb = _MUTED
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.color.rgb = _MUTED
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = _GRID

    # Right column — key stat callouts.
    rows = _drift_rows(pt_hist, last_price)
    active_pts = [r["pt"] for r in rows if r["active"] and r["pt"]]
    consensus_pt = sum(active_pts) / len(active_pts) if active_pts else None
    upside = round((consensus_pt / last_price - 1) * 100, 1) if consensus_pt and last_price else None

    right_left = Inches(8.75)
    right_width = Inches(4.05)
    _stat_callout(slide, right_left, Inches(1.45), right_width, "CONSENSUS PT",
                  f"${consensus_pt:.2f}" if consensus_pt else "—",
                  f"{upside:+.0f}% upside vs last price" if upside is not None else "",
                  _GREEN if (upside or 0) >= 0 else _RED)
    _stat_callout(slide, right_left, Inches(2.55), right_width, "LAST PRICE", f"${last_price:.2f}", "", _ICE)

    if revision_momentum:
        status = revision_momentum.get("status")
        color = _GREEN if status == "GREEN" else _RED if status == "RED" else _ICE
        _stat_callout(slide, right_left, Inches(3.65), right_width, "REVISION MOMENTUM",
                      revision_momentum.get("headline", "—"), revision_momentum.get("detail", ""), color)
        table_top = Inches(5.05)
    else:
        table_top = Inches(3.65)

    # Compact per-firm table — only actively-covering analysts.
    active_rows = [r for r in rows if r["active"]]
    if active_rows:
        n_rows = len(active_rows) + 1
        row_h = Inches(0.32)
        tbl_shape = slide.shapes.add_table(n_rows, 3, right_left, table_top, right_width, row_h * n_rows)
        tbl = tbl_shape.table
        tbl.columns[0].width = Inches(1.85)
        tbl.columns[1].width = Inches(1.05)
        tbl.columns[2].width = Inches(1.15)
        for c, h in enumerate(["Firm", "PT", "8Q Chg"]):
            cell = tbl.cell(0, c)
            cell.text = h
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = _WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = _NAVY_LIGHT
            cell.margin_top = cell.margin_bottom = Pt(2)
        for r_i, r in enumerate(active_rows, start=1):
            chg = r["chg_pct"]
            chg_color = _GREEN if (chg or 0) > 0 else _RED if (chg or 0) < 0 else _MUTED
            vals = [
                (r["firm"], _WHITE), (f"${r['pt']:.2f}" if r["pt"] else "—", _WHITE),
                (f"{chg:+.1f}%" if chg is not None else "—", chg_color),
            ]
            for c, (text, color) in enumerate(vals):
                cell = tbl.cell(r_i, c)
                cell.text = text
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(10)
                run.font.color.rgb = color
                cell.fill.solid()
                cell.fill.fore_color.rgb = _NAVY
                cell.margin_top = cell.margin_bottom = Pt(2)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
